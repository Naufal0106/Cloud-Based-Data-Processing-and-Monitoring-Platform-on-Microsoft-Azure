import sys
import subprocess
from pathlib import Path

# Install python-pptx library if not present
try:
    import pptx
except ImportError:
    print("Installing python-pptx library...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "docs" / "Presentasi_Final_Kelompok_11.pptx"
IMAGE_PATH = ROOT / "docs" / "evidence" / "arsitektur-final.png"

# Color theme
NAVY = RGBColor(15, 76, 129)       # #0f4c81 - Primary Dark
TEAL = RGBColor(15, 118, 110)      # #0f766e - Secondary Accent
DARK_TEXT = RGBColor(30, 41, 59)   # #1e293b - Slate Text
LIGHT_BG = RGBColor(248, 250, 252) # #f8fafc - Light background
WHITE = RGBColor(255, 255, 255)

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title_text, subtitle_text, members):
    # Using blank slide layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, NAVY)

    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle textbox
    txBox_sub = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(1))
    tf_sub = txBox_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = WHITE
    p_sub.alignment = PP_ALIGN.CENTER

    # Members textbox
    txBox_mem = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.3), Inches(2))
    tf_mem = txBox_mem.text_frame
    tf_mem.word_wrap = True
    p_mem = tf_mem.paragraphs[0]
    p_mem.text = "Disusun oleh Kelompok 11:\n" + " | ".join(members)
    p_mem.font.size = Pt(14)
    p_mem.font.color.rgb = WHITE
    p_mem.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title_text, bullets):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, LIGHT_BG)

    # Header Box
    headerBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_h = headerBox.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = title_text
    p_h.font.size = Pt(28)
    p_h.font.bold = True
    p_h.font.color.rgb = NAVY

    # Decorative bottom line for header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    # Content Box
    contentBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
    tf_c = contentBox.text_frame
    tf_c.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p_c = tf_c.paragraphs[0]
        else:
            p_c = tf_c.add_paragraph()
            
        p_c.text = bullet
        p_c.font.size = Pt(15)
        p_c.font.color.rgb = DARK_TEXT
        p_c.space_after = Pt(12)
        
        # Check indentation level (e.g. starts with spaces or sub-points)
        if bullet.strip().startswith("- ") or bullet.strip().startswith("* "):
            p_c.text = bullet.strip()[2:]
            p_c.level = 1
        elif bullet.strip().startswith("• ") or bullet.strip().startswith("  - "):
            p_c.text = bullet.strip()[2:]
            p_c.level = 2
        else:
            p_c.level = 0

def add_image_slide(prs, title_text, img_path):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, LIGHT_BG)

    # Header Box
    headerBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_h = headerBox.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = title_text
    p_h.font.size = Pt(28)
    p_h.font.bold = True
    p_h.font.color.rgb = NAVY

    # Line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    # Image placement
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(1.2), Inches(1.7), width=Inches(10.9), height=Inches(5.0))
    else:
        # Fallback textbox
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.3), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Gambar tidak ditemukan: {img_path.name}]"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

def add_closing_slide(prs, title_text, subtitle_text):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, NAVY)

    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle textbox
    txBox_sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(1))
    tf_sub = txBox_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = WHITE
    p_sub.alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Building final presentation slide deck...")

    # Slide 1: Cover
    add_title_slide(
        prs,
        "LAPORAN TEKNIS AKHIR",
        "Cloud-Based Data Processing & Monitoring Platform on Microsoft Azure",
        [
            "Naufal Ihsan S (DevOps)",
            "M. Arifin Ilham (Backend)",
            "Zhykwa Ceryl M (Architect)",
            "Rendy Saputra (Security)"
        ]
    )

    # Slide 2: Latar Belakang
    add_content_slide(
        prs,
        "Latar Belakang Proyek",
        [
            "Tingginya Kebutuhan Data Real-time: Modul sensor telemetri, logging aplikasi, dan aktivitas pengguna memproduksi data dalam jumlah masif secara kontinu.",
            "Kebutuhan ETL Pipeline yang Handal: Membutuhkan infrastruktur awan yang secara cerdas mampu mengolah berkas besar tanpa lag/hambatan kapasitas.",
            "Tantangan Kualitas Data: Adanya data kosong (missing values) dan data pencilan ekstrem (outliers) di lapangan yang berpotensi merusak hasil keputusan analisis.",
            "Solusi Arsitektur Modern: Kombinasi Hybrid Cloud (Cloudflare Edge CDN + Microsoft Azure) untuk performa unggul serta penghematan biaya operasional secara maksimal."
        ]
    )

    # Slide 3: Tujuan Proyek
    add_content_slide(
        prs,
        "Tujuan Proyek",
        [
            "1. Membangun ETL Pipeline Skalabel: Pipeline pemrosesan data berbasis awan yang dapat memproses file hingga 100 MB secara aman.",
            "2. Pembersihan Data Science Otomatis: Menyediakan proses filtering otomatis untuk anomali dan outliers sebelum disimpan ke database.",
            "3. Penerapan Infrastructure as Code (IaC): Menggunakan Terraform untuk merancang, melacak, dan mereplikasi infrastruktur cloud secara konsisten.",
            "4. Aspek Keamanan & Monitoring Mandiri: Mengedepankan enkripsi transit/at-rest, managed identity, pembatasan port SSH, serta visualisasi metrik observability."
        ]
    )

    # Slide 4: Ringkasan Arsitektur Hybrid
    add_image_slide(prs, "Ringkasan Arsitektur Hybrid", IMAGE_PATH)

    # Slide 5: Infrastruktur Backend & Data Layer
    add_content_slide(
        prs,
        "Infrastruktur Backend & Data Layer",
        [
            "Azure Functions (Serverless Compute):",
            "  - Bertindak sebagai REST API backend menggunakan runtime Python 3.11.",
            "  - Consumption Plan (Y1) memastikan tagihan komputasi dihitung hanya saat kode berjalan.",
            "Azure Blob Storage (Storage Account):",
            "  - Container 'raw-data' menyimpan berkas mentah. Pemicu Blob Trigger otomatis memproses file saat selesai diunggah.",
            "Azure Cosmos DB NoSQL Database:",
            "  - Container 'users' (Partition Key: /email) untuk autentikasi kredensial.",
            "  - Container 'telemetry-data' (Partition Key: /deviceId) untuk data sensor hasil enrichment."
        ]
    )

    # Slide 6: Penerapan Infrastructure as Code (IaC)
    add_content_slide(
        prs,
        "Penerapan Infrastructure as Code (IaC)",
        [
            "Deklarasi Terraform Terstruktur:",
            "  - main.tf & network.tf: Pengaturan provider Azure, resource group, serta setup Virtual Network (VNet) dan Subnet.",
            "  - database.tf & storage.tf: Inisialisasi akun Cosmos DB NoSQL dan Storage Account secara modular.",
            "  - functions.tf & security.tf: Konfigurasi App Service, Function App, Network Security Group (NSG), dan Azure Key Vault.",
            "  - monitoring.tf: Setup Application Insights dan Rules Alert otomatis.",
            "Keuntungan Utama:",
            "  - Mencegah inkonsistensi setelan antara tahapan dev, staging, dan prod.",
            "  - Mempermudah audit infrastruktur langsung dari repository Git."
        ]
    )

    # Slide 7: Keamanan & Pengolahan Secret Kredensial
    add_content_slide(
        prs,
        "Keamanan & Pengolahan Secret Kredensial",
        [
            "Isolasi Secret di Azure Key Vault:",
            "  - Database connection string, token API, dan kunci rahasia JWT dipisahkan sepenuhnya dari repository kode program.",
            "Integrasi Managed Identity (Passwordless):",
            "  - Azure Functions mengakses Key Vault melalui System-Assigned Managed Identity.",
            "  - Tidak ada kredensial yang ditulis secara manual di file konfigurasi (mengurangi resiko kebocoran akses).",
            "Same-Origin Proxy pada Edge:",
            "  - Frontend memanggil path '/api/*', lalu Cloudflare Pages Function menyisipkan API Key rahasia di sisi server CDN sebelum meneruskan ke Azure."
        ]
    )

    # Slide 8: Hardening Akses Port SSH 22
    add_content_slide(
        prs,
        "Hardening Akses Port SSH 22",
        [
            "Identifikasi Resiko Keamanan VM:",
            "  - Pada setelan awal, port SSH 22 pada VM eksplorasi terbuka bebas untuk lalu lintas IP publik internet.",
            "Langkah Mitigasi Hardening di Terraform (security.tf):",
            "  - Mengonfigurasi rule inbound NSG (Network Security Group).",
            "  - Membatasi sumber akses port SSH 22 secara eksklusif hanya untuk IP publik admin terdaftar via variabel 'var.admin_ssh_allowed_ip'.",
            "Hasil Proteksi VM:",
            "  - Menghalangi potensi brute-force SSH dan scanning port berbahaya dari peretas luar."
        ]
    )

    # Slide 9: Logika Pembersihan Data Science
    add_content_slide(
        prs,
        "Logika Pembersihan Data Science",
        [
            "Alur Pembersihan Otomatis di Backend:",
            "  1. Pembersihan Nilai Kosong (Missing Value Check):",
            "     - Baris yang memiliki cell kosong pada kolom krusial akan dilewati/dihapus agar tidak mendistorsi statistik rata-rata.",
            "  2. Pembersihan Nilai Pencilan (Outlier IQR Filter):",
            "     - Menggunakan rumus statistik Interquartile Range (IQR):",
            "       • Rentang IQR = Q3 - Q1",
            "       • Batas Bawah = Q1 - 1.5 * IQR",
            "       • Batas Atas = Q3 + 1.5 * IQR",
            "     - Baris dengan nilai di luar batas tersebut disaring sebelum data disimpan ke Cosmos DB.",
            "  3. Visualisasi Bersih:",
            "     - Dashboard menyajikan visualisasi data yang sudah bebas anomali & outliers."
        ]
    )

    # Slide 10: Pengujian Fungsional End-to-End
    add_content_slide(
        prs,
        "Pengujian Fungsional End-to-End",
        [
            "5 Skenario Uji Sukses Utama:",
            "  1. Health Check (GET /api/hello): Validasi ketersediaan serverless API.",
            "  2. User Registration (POST /api/register): Berhasil menyimpan hash sandi ke Cosmos DB users.",
            "  3. User Authentication (POST /api/login): Menerbitkan JWT Token tepercaya untuk sesi 8 jam.",
            "  4. Data Upload (POST /api/upload): Pengunggahan berkas besar dengan status processing yang lancar.",
            "  5. Role-Based Access Control (RBAC): Memastikan akun biasa diblokir saat mencoba mengakses halaman konfigurasi admin.",
            "Alat Pengujian:",
            "  - Otomatisasi script tes via Windows PowerShell (scripts/test-auth-db.ps1)."
        ]
    )

    # Slide 11: Observability & Alerting
    add_content_slide(
        prs,
        "Observability & Alerting",
        [
            "Centralized Logging & Monitoring:",
            "  - Metrik latency, working set memory, exception, dan HTTP error code dicatat oleh Application Insights secara terpusat.",
            "Azure Monitor Alert Rules:",
            "  - 1. Alert 5xx Error: Mengirim notifikasi email otomatis ke tim operasi jika backend menghasilkan status HTTP 5xx > 5 kali dalam 5 menit.",
            "  - 2. Alert Latency: Peringatan aktif jika respon rata-rata backend API > 2 detik.",
            "  - 3. Alert VM CPU: Peringatan jika utilisasi CPU VM admin > 80% dalam 5 menit.",
            "Action Group:",
            "  - Menghubungkan log deteksi insiden langsung ke kotak masuk email tim pengawas."
        ]
    )

    # Slide 12: Analisis & Optimasi Biaya
    add_content_slide(
        prs,
        "Analisis & Optimasi Biaya",
        [
            "Fixed Cost (CapEx):",
            "  - Pembelian nama domain kustom 'kelompok11cc.my.id' melalui registrar lokal Domainesia dengan biaya Rp 15.000,- / tahun.",
            "Operational Cost (OpEx) & Optimasi:",
            "  - Azure Functions Consumption Mode: Menghindari biaya server idle (hanya bayar saat program dieksekusi).",
            "  - Cosmos DB Serverless: Menghilangkan kapasitas provisioned RU/s yang tidak terpakai.",
            "  - VM Scheduler: Mematikan VM di luar kebutuhan demo.",
            "  - Storage Lifecycle Policy: Otomatisasi pembersihan berkas mentah berumur tua di Blob Storage untuk menekan biaya penyimpanan."
        ]
    )

    # Slide 13: Demo Sistem (Live Show)
    add_content_slide(
        prs,
        "Demo Sistem (Live Show)",
        [
            "Skenario Demonstrasi:",
            "  - 1. Registrasi Akun & Login User Baru.",
            "  - 2. Unggah Data Tanpa Pembersihan: Tunjukkan visual data mentah beserta peringatan missing value & outliers.",
            "  - 3. Unggah Data Dengan Pembersihan: Tunjukkan grafik histogram bersih dan quality score yang melonjak naik.",
            "  - 4. Panel Admin: Tunjukkan antarmuka pengelolaan peran anggota tim.",
            "  - 5. Panel Developer: Tunjukkan grafik telemetri performa backend (CPU, Request Rate, Latency)."
        ]
    )

    # Slide 14: Keterbatasan & Saran Pengembangan
    add_content_slide(
        prs,
        "Keterbatasan & Saran Pengembangan",
        [
            "Keterbatasan Sistem Saat Ini:",
            "  - Efek Cold Start pada komputasi serverless yang memerlukan jeda beberapa detik pada pemanggilan API pertama.",
            "  - Batas maksimal unggahan berkas tunggal dibatasi pada ukuran 100 MB.",
            "Rekomendasi Pengembangan Selanjutnya:",
            "  - Mengintegrasikan Azure Event Hubs untuk pemrosesan data streaming kontinu dari sensor IoT.",
            "  - Menggunakan modul Machine Learning tertanam untuk mendeteksi pencilan data secara prediktif.",
            "  - Penerapan Azure VNet Integration untuk memisahkan lalu lintas database dari internet publik secara total."
        ]
    )

    # Slide 15: Penutup
    add_closing_slide(prs, "TERIMA KASIH", "Sesi Tanya Jawab Dibuka | Kelompok 11 UPR")

    # Save
    prs.save(str(PPTX_PATH))
    print(f"Presentation saved successfully to {PPTX_PATH}!")

if __name__ == "__main__":
    main()
